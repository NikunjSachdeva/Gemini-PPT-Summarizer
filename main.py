import streamlit as st
import google.generativeai as genai
from pptx import Presentation
import os

# Configure Google Gemini API
genai.configure(api_key="AIzaSyB6uQooQ4Bp-_vo1uj2PVJ9ppbTvqoYX3I")
model = genai.GenerativeModel("gemini-pro")

# Function to extract text from PowerPoint
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    text = ""
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n\n"
    
    return text

# Function to summarize content
def summarize_text(text):
    prompt = f"Summarize the following presentation text and list key topics:\n\n{text}"
    response = model.generate_content(prompt)
    return response.text

# Function for question-answering
def answer_question(text, question):
    prompt = f"The following text is extracted from a PowerPoint presentation. Answer the question based on it:\n\n{text}\n\nQuestion: {question}"
    response = model.generate_content(prompt)
    return response.text

# Streamlit UI
st.title("📊 PPTX AI Summarizer & Q&A Assistant")

uploaded_file = st.file_uploader("Upload a PowerPoint file (.pptx)", type=["pptx"])

if uploaded_file is not None:
    st.success("✅ File uploaded successfully!")
    
    # Save uploaded file temporarily
    ppt_path = "uploaded_presentation.pptx"
    with open(ppt_path, "wb") as f:
        f.write(uploaded_file.read())

    # Extract text
    ppt_text = extract_text_from_ppt(ppt_path)
    
    # Remove temporary file
    os.remove(ppt_path)
    
    if ppt_text.strip():
        st.subheader("📜 Extracted Text (Hidden by Default)")
        with st.expander("Click to view extracted text"):
            st.text_area("Extracted Text", ppt_text, height=300)

        # Summarize the presentation
        if st.button("🔍 Summarize"):
            with st.spinner("Generating summary..."):
                summary = summarize_text(ppt_text)
            st.subheader("📖 Summary & Topics")
            st.write(summary)

        # Q&A Section
        st.subheader("💡 Ask a Question")
        user_question = st.text_input("Enter your question:")
        if st.button("Ask Gemini"):
            if user_question.strip():
                with st.spinner("Thinking..."):
                    answer = answer_question(ppt_text, user_question)
                st.subheader("🤖 AI Response")
                st.write(answer)
            else:
                st.warning("⚠️ Please enter a question.")
    else:
        st.error("⚠️ No text found in the presentation.")
